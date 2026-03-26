import io as _io
import base64, json, os, requests, streamlit as st, time, sys
from datetime import datetime, timedelta
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Add parent directory to path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ollama_helper import check_ollama_connection, generate_story_stream, get_model_name, check_guardrails
from db import (
    init_db, create_user, verify_user, save_story,
    get_stories, get_story_by_id, delete_story, get_user_info, format_date,
)

st.set_page_config(
    page_title="Story Waver",
    page_icon="📖",
    layout="wide",
    initial_sidebar_state="expanded",
)

init_db()

# ── CSS ───────────────────────────────────────────────────────────────
st.markdown("""
<style>
/* ── Global ── */
.stApp { background-color: #fff5f7 !important; color: #1e3a8a; }
[data-testid="stSidebar"] { background-color: #ffe4e6; border-right: 1px solid #fbc4c9; }
[data-testid="stSidebar"] * { color: #1e3a8a !important; }
/* Fix for sidebar genre selection - remove dark box */
[data-testid="stSidebar"] .stSelectbox > div {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 8px !important;
    color: #1e3a8a !important;
}

[data-testid="stSidebar"] .stSelectbox > div:hover {
    border-color: #ff69b4 !important;
}

[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] {
    background: transparent !important;
    border: none !important;
}

[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] span {
    color: #1e3a8a !important;
    background: transparent !important;
}

/* Fix for all dropdown menus in sidebar */
[data-testid="stSidebar"] .stSelectbox div[data-baseweb="select"] > div {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: none !important;
}

/* Fix for the dropdown arrow */
[data-testid="stSidebar"] .stSelectbox svg {
    fill: #ff69b4 !important;
    color: #ff69b4 !important;
}

/* Fix for selectbox labels */
[data-testid="stSidebar"] .stSelectbox label {
    color: #1e3a8a !important;
    font-weight: 500 !important;
    font-size: 0.75rem !important;
}

/* Fix for multi-select if present */
[data-testid="stSidebar"] .stMultiSelect > div {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 8px !important;
}

/* Fix for slider */
[data-testid="stSidebar"] .stSlider > div {
    background: transparent !important;
}

[data-testid="stSidebar"] .stSlider div[data-baseweb="slider"] {
    background: transparent !important;
}

/* Fix for number input */
[data-testid="stSidebar"] .stNumberInput > div {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 8px !important;
}

/* Ensure all sidebar containers have consistent styling */
[data-testid="stSidebar"] .stVerticalBlock {
    background: transparent !important;
    border: none !important;
    box-shadow: none !important;
    padding: 0 !important;
    margin: 0 !important;
}

/* Fix for any remaining dark elements */
[data-testid="stSidebar"] [data-testid="baseButton-secondary"] {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    color: #1e3a8a !important;
}

[data-testid="stSidebar"] [data-testid="baseButton-secondary"]:hover {
    border-color: #ff69b4 !important;
}

/* Fix for all expanders globally (Search and Filter, etc.) */
[data-testid="stExpander"] details,
[data-testid="stExpander"] details summary,
[data-testid="stExpander"] details > div {
    background: linear-gradient(145deg, #ff99cc, #ff69b4) !important;
}
[data-testid="stExpander"] {
    background: transparent !important;
    border: 1px solid #ff1493 !important;
    border-radius: 8px !important;
    margin-bottom: 8px !important;
    overflow: hidden !important;
}
[data-testid="stExpander"] summary, [data-testid="stExpander"] summary * {
    color: #1e3a8a !important;
    font-weight: 600 !important;
    background-color: transparent !important;
}

/* Fix for any dropdown options */
[data-testid="stSidebar"] [data-baseweb="menu"] {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 8px !important;
}

[data-testid="stSidebar"] [data-baseweb="menu"] li {
    color: #1e3a8a !important;
    background: transparent !important;
}

[data-testid="stSidebar"] [data-baseweb="menu"] li:hover {
    background: #ffb6c1 !important;
    color: #1e3a8a !important;
}

/* Fix for the bottom sidebar container */
[data-testid="stVerticalBlock"]:has(.sidebar-section-header) {
    background: linear-gradient(145deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 12px !important;
    padding: 18px 24px !important;
    margin: 16px 0 !important;
    box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1) !important;
}

/* Ensure all text in sidebar is navy blue */
[data-testid="stSidebar"] * {
    color: #1e3a8a !important;
}

/* Fix for the genre tag inside the main content */
.msg-tag.tag-genre {
    background: #ff69b4 !important;
    color: white !important;
}

/* Global input field styling (Search, Login, etc.) */
.stTextInput div[data-baseweb="input"],
.stTextArea div[data-baseweb="textarea"],
.stNumberInput div[data-baseweb="input"] {
    background: linear-gradient(145deg, #ff99cc, #ff69b4) !important;
    border: 1px solid #ff1493 !important;
    border-radius: 8px !important;
}
/* Pierce through dark mode sub-wrappers */
.stTextInput div[data-baseweb="input"] *,
.stTextArea div[data-baseweb="textarea"] *,
.stNumberInput div[data-baseweb="input"] * {
    background-color: transparent !important;
    background: transparent !important;
    color: #1e3a8a !important;
}

/* Ensure the sidebar section headers remain pink gradient */
.sidebar-section-header {
    background: linear-gradient(90deg, #ff69b4, #ff1493) !important;
    color: white !important;
    padding: 8px 12px !important;
    border-radius: 6px !important;
    margin: 16px 0 8px !important;
    font-size: 0.85rem !important;
    font-weight: 700 !important;
    text-transform: uppercase !important;
    letter-spacing: 1px !important;
}

/* Fix for checkboxes in the bottom sidebar */
[data-testid="stSidebar"] .stCheckbox {
    background: transparent !important;
}

[data-testid="stSidebar"] .stCheckbox label {
    color: #1e3a8a !important;
}

/* ── Title ── */
.main-title {
    text-align: center; font-size: 2.6rem; font-weight: 900; letter-spacing: 8px;
    background: linear-gradient(135deg, #ff69b4 0%, #ffb6c1 50%, #ff1493 100%);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text; margin-bottom: 0; padding-top: 10px;
}
.main-subtitle {
    text-align: center; color: #ff1493; font-size: 0.9rem;
    margin-top: 2px; margin-bottom: 10px; letter-spacing: 1px;
}

/* ── Chat messages ── */
.stChatMessageContainer {
    background: transparent !important;
    padding: 0 !important;
}

[data-testid="stChatMessageContent"] {
    background: linear-gradient(145deg, #fff0f5, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 12px !important;
    padding: 18px 22px !important;
    color: #1e3a8a !important;
    line-height: 1.85 !important;
    font-size: 1.0rem !important;
}
/* User bubble — slightly different tint */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) [data-testid="stChatMessageContent"] {
    background: linear-gradient(145deg, #fff0f5, #ffe4ec) !important;
    border-color: #ffb6c1 !important;
    color: #1e3a8a !important;
}

/* ── Chat input bar ── */
[data-testid="stChatInput"] {
    background: linear-gradient(145deg, #fff0f5, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 14px !important;
}
[data-testid="stChatInput"] textarea {
    background: transparent !important;
    color: #1e3a8a !important;
    font-size: 0.97rem !important;
    font-weight: 600 !important;
}
[data-testid="stChatInput"] textarea::placeholder {
    color: #666666 !important;
    opacity: 1 !important;
    font-weight: 500 !important;
}
[data-testid="stChatInput"] button {
    background: linear-gradient(135deg, #ff69b4, #ff1493) !important;
    border-radius: 8px !important;
}
/* Bottom padding so the last message isn't hidden behind the input */
.main > div { padding-bottom: 90px !important; }

/* ── Genre / Tone tags ── */
.msg-tag {
    display: inline-block; font-size: 0.66rem; font-weight: 700;
    letter-spacing: 1px; text-transform: uppercase;
    padding: 2px 8px; border-radius: 20px; margin-right: 5px; margin-bottom: 8px;
}
.tag-genre { 
    background: #ff69b4; 
    color: white !important;
}
.tag-tone  { 
    background: #ffe4ec; 
    color: #1e3a8a !important; 
    border: 1px solid #ffb6c1; 
}
.tag-mode  { 
    background: #ffd9e6; 
    color: #1e3a8a !important; 
    border: 1px solid #ffb6c1; 
}

/* Selected tag styles - white text for selected items */
.tag-genre.selected, 
.tag-tone.selected, 
.tag-mode.selected {
    background: #ff1493 !important;
    color: white !important;
    box-shadow: 0 2px 8px rgba(255, 20, 147, 0.3);
}

/* ── Sidebar section header ── */
.sidebar-section-header {
    background: linear-gradient(90deg, #ff69b4, #ff1493);
    color: white !important;
}

/* ── Plot steering section ── */
.plot-steering-select {
    background: linear-gradient(145deg, #ffe4ec, #ffd9e6) !important; border: 1px solid #ffb6c1 !important;
    border-radius: 8px !important; color: #1e3a8a !important;
}
.plot-steering-select[data-baseweb="select"] span {
    color: #1e3a8a !important;
}

/* ── Collaborative mode section ── */
.collaborative-indicator {
    background: linear-gradient(145deg, #ffffff, #ffe4ec);
    border: 1px solid #ffb6c1; border-radius: 10px;
    padding: 12px; margin: 8px 0;
}
.collaborative-turn {
    color: #ff1493; font-weight: 600; font-size: 0.8rem;
}
.collaborative-user {
    color: #1e3a8a; font-size: 0.9rem;
}
.collaborative-info {
    color: #666; font-size: 0.75rem; margin-top: 4px;
}

/* ── User badge ── */
.user-badge {
    background: linear-gradient(135deg, #ff69b4, #ff1493) !important;
    color: white !important;
    padding: 4px 12px !important;
    border-radius: 20px !important;
    font-size: 0.85rem !important;
    font-weight: 600 !important;
    display: inline-block !important;
    margin: 4px 0 !important;
}

/* ── Status indicators ── */
.status-online {
    color: #28a745 !important;
    font-weight: 600 !important;
    font-size: 0.8rem !important;
}
.status-offline {
    color: #dc3545 !important;
    font-weight: 600 !important;
    font-size: 0.8rem !important;
}

/* ── Sidebar history card ── */
.history-card {
    background: linear-gradient(145deg, #ffe4ec, #ffd9e6) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 8px !important;
    padding: 10px !important;
    margin-bottom: 8px !important;
}
.history-prompt-label {
    font-size: 0.7rem !important;
    color: #ff69b4 !important;
    font-weight: 600 !important;
    text-transform: uppercase !important;
    margin-bottom: 2px !important;
}
.history-title {
    font-size: 0.8rem !important;
    color: #1e3a8a !important;
    margin-bottom: 4px !important;
    line-height: 1.3 !important;
}
.history-date {
    font-size: 0.7rem !important;
    color: #ff69b4 !important;
}

/* ── Stats dashboard ── */
.stats-card {
    background: linear-gradient(135deg, #ffffff, #ffe4ec) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 12px !important;
    padding: 16px !important;
    margin-bottom: 12px !important;
    text-align: center !important;
}
.stats-number {
    font-size: 1.8rem !important;
    font-weight: 700 !important;
    color: #ff69b4 !important;
    margin-bottom: 4px !important;
}
.stats-label {
    font-size: 0.8rem !important;
    color: #666 !important;
    text-transform: uppercase !important;
    letter-spacing: 1px !important;
}
.stButton > button[kind="primary"]:hover {
    background: linear-gradient(135deg, #ff1493, #ff69b4);
    box-shadow: 0 4px 20px rgba(255, 105, 180, 0.4);
}

/* ── Status ── */
.status-online  { color: #28a745; font-weight: 700; font-size: 0.76rem; }
.status-offline { color: #dc3545; font-weight: 700; font-size: 0.76rem; }

/* ── Input fields ── */
.stTextArea textarea, .stTextInput input {
    background: linear-gradient(145deg, #ffe4ec, #ffd9e6) !important; color: #1e3a8a !important;
    border: 1px solid #ffb6c1 !important; border-radius: 8px !important;
}
/* Specific styling for the prompt text area */
.stTextArea textarea[placeholder*="Describe your story here"] {
    color: #ff69b4 !important;
}
.stTextArea textarea::placeholder {
    color: #ff69b4 !important;
}
.stTextArea textarea:focus {
    color: #ff69b4 !important;
}
.stTextArea textarea {
    color: #ff69b4 !important;
}
/* Main chat input styling */
[data-testid="stChatInput"] textarea {
    color: #ff69b4 !important;
}
[data-testid="stChatInput"] textarea::placeholder {
    color: #ff69b4 !important;
}
[data-testid="stChatInput"] textarea:focus {
    color: #ff69b4 !important;
}

/* ── History page ── */
.story-view-header {
    background: linear-gradient(90deg, #ff69b4, #ff1493);
    border: 1px solid #ffb6c1; border-radius: 10px;
    padding: 14px 18px; margin-bottom: 16px;
    color: white !important;
}
.story-view-prompt { color: #ff69b4 !important; font-size: 0.97rem; margin-top: 8px; line-height: 1.5; }
.story-view-prompt span { color: #ff1493 !important; font-weight: 600; margin-right: 4px; }
.story-view-meta { color: #ff69b4 !important; font-size: 0.78rem; margin-top: 6px; }
.story-box {
    background: linear-gradient(145deg, #fff0f5, #ffe4ec) !important;
    border: 1px solid #ffb6c1; border-radius: 12px;
    padding: 24px 28px; line-height: 1.85; font-size: 1rem;
    color: #1e3a8a !important;
    white-space: pre-wrap;
}

hr { border-color: #ffb6c1 !important; }
[data-testid="stExpander"] { background: #ffffff; border: 1px solid #ffb6c1 !important; border-radius: 8px; }

/* ── Stats Row ── */
.stats-row { display: flex; gap: 12px; margin: 0 0 24px; }
.stat-card {
    flex: 1; background: #ffffff; border: 1px solid #ffb6c1;
    border-radius: 14px; padding: 18px 12px; text-align: center;
    transition: border-color 0.3s, transform 0.2s;
}
.stat-card:hover { border-color: #ff69b4; transform: translateY(-2px); }
.stat-icon { font-size: 1.4rem; display: block; margin-bottom: 8px; }
.stat-value {
    font-size: 1.9rem; font-weight: 800; display: block;
    background: linear-gradient(135deg, #ff69b4, #ff1493);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text; line-height: 1.2; margin-bottom: 5px;
}
.stat-label { color: #ff69b4 !important; font-size: 0.63rem; font-weight: 700; letter-spacing: 1.8px; text-transform: uppercase; display: block; }

/* ── Story Progress Bar ── */
.progress-wrap {
    background: #ffffff; border: 1px solid #ffb6c1; border-radius: 14px;
    padding: 18px 24px; margin-bottom: 24px;
}
.progress-title { color: #1e3a8a !important; font-size: 0.68rem; font-weight: 700; letter-spacing: 2px; text-transform: uppercase; margin-bottom: 14px; }
.progress-phases { display: flex; align-items: center; }
.phase-node {
    display: flex; flex-direction: column; align-items: center; gap: 6px; flex: 0 0 auto;
}
.phase-dot {
    width: 12px; height: 12px; border-radius: 50%;
    background: #ffe4ec; border: 2px solid #ffb6c1;
    transition: all 0.3s;
}
.phase-dot.active { background: #ff69b4; border-color: #ff69b4; box-shadow: 0 0 10px rgba(255, 105, 180, 0.4); }
.phase-dot.done   { background: #ffb6c1; border-color: #ff69b4; }
.phase-label { font-size: 0.68rem; color: #ff69b4 !important; font-weight: 600; letter-spacing: 0.5px; }
.phase-label.active { color: #ff1493 !important; }
.phase-label.done   { color: #ff69b4 !important; }
.phase-line { flex: 1; height: 2px; background: #ffe4ec; margin: 0 6px; margin-bottom: 18px; border-radius: 2px; }
.phase-line.done { background: linear-gradient(90deg, #ff69b4, #ffb6c1); }
.progress-words { color: #ff69b4 !important; font-size: 0.72rem; margin-top: 12px; font-weight: 500; }

/* ── Bottom Sidebar Styling ── */
[data-testid="stVerticalBlock"] {
    background: linear-gradient(180deg, #ffe4ec 0%, #ffd9e6 100%) !important;
    border: 1px solid #ffb6c1 !important;
    border-radius: 12px !important;
    padding: 18px 24px !important;
    margin: 16px 0 !important;
    box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1) !important;
}
[data-testid="stVerticalBlock"] * {
    color: #1e3a8a !important;
    text-shadow: none !important;
}
[data-testid="stVerticalBlock"] label {
    color: #1e3a8a !important;
    font-weight: 500 !important;
    font-size: 0.75rem !important;
    text-align: left !important;
    display: block !important;
    margin-bottom: 4px !important;
}
[data-testid="stVerticalBlock"] small {
    color: #1e3a8a !important;
    font-weight: 500 !important;
    font-size: 0.7rem !important;
    text-align: left !important;
    display: block !important;
}
[data-testid="stVerticalBlock"] span {
    color: #1e3a8a !important;
    font-weight: 500 !important;
    font-size: 0.75rem !important;
    text-align: left !important;
}
[data-testid="stVerticalBlock"] .collaborative-info {
    color: #ff69b4 !important;
    font-size: 0.75rem !important;
    margin-top: 4px !important;
}

/* ── Genre Background Overlay ── */
.genre-overlay {
    position: fixed; inset: 0; pointer-events: none; z-index: 1; overflow: hidden;
}
.gb-elem {
    position: absolute; bottom: -60px; opacity: 0; user-select: none;
    animation-iteration-count: infinite; animation-timing-function: linear;
}
@keyframes gbFloatUp {
    0%   { transform: translateY(0) rotate(0deg);       opacity: 0; }
    8%   { opacity: var(--gb-opacity, 0.15); }
    92%  { opacity: var(--gb-opacity, 0.15); }
    100% { transform: translateY(-110vh) rotate(25deg);  opacity: 0; }
}
@keyframes gbDrift {
    0%   { transform: translate(0,0) rotate(0deg);           opacity: 0; }
    8%   { opacity: var(--gb-opacity, 0.15); }
    40%  { transform: translate(40px,-45vh) rotate(180deg); opacity: var(--gb-opacity, 0.15); }
    60%  { transform: translate(-30px,-65vh) rotate(270deg); opacity: var(--gb-opacity, 0.15); }
    92%  { transform: translate(10px,-100vh) rotate(340deg); opacity: var(--gb-opacity, 0.15); }
    100% { transform: translate(10px,-110vh) rotate(360deg); opacity: 0; }
}
@keyframes gbTwinkle {
    0%,100% { opacity: 0;                      transform: scale(0.5) translateY(0);    }
    30%,70% { opacity: var(--gb-opacity, 0.2); transform: scale(1.1) translateY(-10px); }
}
@keyframes gbBounce {
    0%   { transform: translateY(0) scale(1) rotate(0deg);    opacity: 0; }
    5%   { opacity: var(--gb-opacity, 0.15); }
    20%  { transform: translateY(-20vh) scale(1.1) rotate(10deg); }
    50%  { transform: translateY(-55vh) scale(0.9) rotate(-10deg); }
    80%  { transform: translateY(-85vh) scale(1.05) rotate(5deg); }
    95%  { opacity: var(--gb-opacity, 0.15); }
    100% { transform: translateY(-110vh) scale(0.8) rotate(0deg); opacity: 0; }
}
@keyframes gbPulse {
    0%,100% { opacity: 0.03; transform: scale(0.9); }
    50%     { opacity: var(--gb-opacity, 0.18); transform: scale(1.15) translateY(-5px); }
}
@keyframes gbSway {
    0%   { transform: translateY(0) rotate(-8deg);     opacity: 0; }
    8%   { opacity: var(--gb-opacity, 0.15); }
    25%  { transform: translateY(-28vh) rotate(8deg);  opacity: var(--gb-opacity, 0.15); }
    50%  { transform: translateY(-55vh) rotate(-8deg);  opacity: var(--gb-opacity, 0.15); }
    75%  { transform: translateY(-80vh) rotate(8deg);  opacity: var(--gb-opacity, 0.15); }
    92%  { opacity: var(--gb-opacity, 0.15); }
    100% { transform: translateY(-110vh) rotate(-8deg); opacity: 0; }
}

/* ── Auth Page Full Screen Professional Styling ── */
.auth-fullscreen {
    min-height: 100vh;
    display: flex;
    align-items: center;
    justify-content: center;
    background: linear-gradient(135deg, #fff5f7 0%, #ffe4ec 50%, #ffd9e6 100%);
    padding: 20px;
}
.auth-content {
    width: 100%;
    max-width: 500px;
    text-align: center;
}
.auth-title {
    font-size: 4rem;
    font-weight: 900;
    letter-spacing: 12px;
    background: linear-gradient(135deg, #ff69b4, #ff1493);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 10px;
    line-height: 1.2;
}
.auth-subtitle {
    font-size: 1.3rem;
    color: #ff1493;
    margin-bottom: 30px;
    font-weight: 500;
    letter-spacing: 2px;
}
/* Make tabs container full width */
.stTabs [data-baseweb="tab-list"] {
    gap: 12px;
    background: #fff0f5;
    border-radius: 16px;
    padding: 6px;
    margin-bottom: 30px;
    border: 1px solid #ffb6c1;
    display: flex;
    justify-content: center;
    align-items: center;
    max-width: 600px;
    margin-left: auto;
    margin-right: auto;
}
.stTabs [data-baseweb="tab"] {
    border-radius: 12px;
    padding: 12px 24px;
    color: #1e3a8a !important;
    font-weight: 600;
    font-size: 1rem;
    transition: all 0.3s;
}
.stTabs [data-baseweb="tab"][aria-selected="true"] {
    background: linear-gradient(135deg, #ff69b4, #ff1493);
    color: white !important;
}
/* Style the forms to be centered but not constrained */
.stForm {
    max-width: 400px;
    margin-left: auto;
    margin-right: auto;
    padding: 20px;
    background: white;
    border-radius: 24px;
    box-shadow: 0 20px 60px rgba(255, 105, 180, 0.25);
    border: 1px solid rgba(255, 182, 193, 0.3);
}
.stTextInput label {
    color: #1e3a8a !important;
    font-weight: 600 !important;
    font-size: 0.95rem !important;
    margin-bottom: 8px !important;
}
.stTextInput div[data-baseweb="input"] {
    background: linear-gradient(145deg, #ff99cc, #ff69b4) !important;
    border: 2px solid #ff1493 !important;
    border-radius: 12px !important;
}
.stTextInput input {
    background: transparent !important;
    padding: 14px 16px !important;
    color: #1e3a8a !important;
    font-size: 1rem !important;
    transition: all 0.3s !important;
    caret-color: #ff1493 !important;
    text-align: left !important;
}
[data-testid="InputInstructions"], .stTextInput small, .stForm small, div[data-testid="stFormSubmitInstructions"] {
    display: none !important;
}
.stTextInput input:focus {
    border-color: #ff69b4 !important;
    box-shadow: 0 0 0 4px rgba(255, 105, 180, 0.15) !important;
    outline: none !important;
}
.stTextInput input::placeholder {
    color: #ffb6c1 !important;
    opacity: 0.8 !important;
}
.stTextInput input:hover {
    border-color: #ff69b4 !important;
}
.stButton > button {
    background: linear-gradient(135deg, #ff69b4, #ff1493) !important;
    color: white !important;
    border: none !important;
    border-radius: 12px !important;
    padding: 14px 24px !important;
    font-weight: 600 !important;
    font-size: 1rem !important;
    transition: all 0.3s !important;
    border: 2px solid transparent !important;
    width: 100%;
    margin-top: 10px;
}
.stButton > button:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 25px rgba(255, 105, 180, 0.4) !important;
    border-color: #ffb6c1 !important;
}
.stAlert {
    border-radius: 12px !important;
    border-left: 4px solid #ff69b4 !important;
    background: #fff5f7 !important;
    color: #1e3a8a !important;
    margin-top: 20px !important;
    max-width: 400px;
    margin-left: auto;
    margin-right: auto;
}
.stAlert > div {
    color: #1e3a8a !important;
}
/* Ensure cursor is visible in all browsers */
input, textarea {
    caret-color: #ff1493 !important;
}
/* Add some breathing room at the bottom */
.main > div {
    padding-bottom: 40px !important;
}

/* ── Custom UI Fixes (Iteration 4) ── */
/* 1. Fix Genre and all Selectboxes dropdown menus globally - darker soft pink */
.stSelectbox > div { 
    background: linear-gradient(145deg, #ff99cc, #ff69b4) !important; 
    border: 1px solid #ff1493 !important; 
    border-radius: 8px !important; 
}
.stSelectbox div[data-baseweb="select"] > div { 
    background: transparent !important; 
    border: none !important; 
}
.stSelectbox [data-baseweb="select"] span { 
    color: #1e3a8a !important; 
    font-weight: 500 !important;
    background: transparent !important; 
}
.stSelectbox svg { 
    fill: #1e3a8a !important; 
    color: #1e3a8a !important; 
}
[data-baseweb="menu"], [data-baseweb="popover"], div[data-testid="stSelectbox"] div[role="listbox"] { 
    background: linear-gradient(145deg, #ff99cc, #ff69b4) !important; 
    border: 1px solid #ff1493 !important; 
    border-radius: 8px !important; 
}
[data-baseweb="menu"] li, div[role="option"] { 
    color: #1e3a8a !important; 
    font-weight: 500 !important;
    background: transparent !important; 
}
[data-baseweb="menu"] li:hover, div[role="option"]:hover { 
    background: #ff1493 !important; 
    color: #ffffff !important; 
}

/* 2. Seamless 'Describe Your Story Idea' Section (Force Black Box removal) */
[data-testid="stBottom"], 
[data-testid="stBottom"] > div,
[data-testid="stBottomBlockContainer"],
[data-testid="stChatInputContainer"],
.stChatFloatingInputContainer {
    background: transparent !important;
    background-color: transparent !important;
    border: none !important;
    box-shadow: none !important;
}
/* Ensure entire background stack inside ChatInput is clear before styling the main box */
[data-testid="stChatInput"], [data-testid="stChatInput"] * {
    background-color: transparent !important;
    border: none !important;
}
[data-testid="stChatInput"] {
    background: linear-gradient(135deg, #ff99cc, #ff69b4) !important;
    border: 1px solid #ff1493 !important; 
    box-shadow: 0 4px 15px rgba(255, 105, 180, 0.2) !important;
    border-radius: 14px !important;
    padding: 0 !important;
}
[data-testid="stChatInput"] textarea {
    color: #1e3a8a !important;
    font-weight: 500 !important;
    background: transparent !important;
    background-color: transparent !important;
}
[data-testid="stChatInput"] textarea::placeholder {
    color: #1e3a8a !important;
    opacity: 0.7 !important;
}
[data-testid="stChatInput"] button {
    background: #ff1493 !important;
    border-radius: 8px !important;
    margin: 4px !important;
}

/* 3. Search and Filter From / To inputs - pure background coloring without breaking layout */
div[data-testid="stDateInput"] > label {
    margin-bottom: 2px !important;
    color: #1e3a8a !important;
    font-size: 0.9rem !important;
    font-weight: 600 !important;
}
.stDateInput div[data-baseweb="input"] {
    background: linear-gradient(145deg, #ff99cc, #ff69b4) !important;
    border: 1px solid #ff1493 !important;
    border-radius: 8px !important;
}
/* Force transparency recursively */
.stDateInput div[data-baseweb="input"] * {
    background-color: transparent !important;
    background: transparent !important;
}
.stDateInput input {
    color: #1e3a8a !important;
    font-weight: 600 !important;
    text-align: center !important;
    padding: 6px 4px !important; 
    font-size: 0.85rem !important; 
}
</style>
""", unsafe_allow_html=True)

# ── Session State ─────────────────────────────────────────────────────
for key, default in {
    "logged_in": False,
    "username": "",
    "messages": [],           # chat turns: {role, content, genre, tone, writing_mode}
    "view_story_id": None,
    "page": "main",
    "collaborative_mode": False,
    "collaborative_turn": 1,  # Track whose turn it is (1 or 2)
    "collaborative_contributions": [],  # Store user contributions separately
    "selected_story_elements": [],  # Store selected story elements
    "selected_story_theme": "",  # Store selected story theme
    "story_characters": {},  # Store user-defined characters
}.items():
    if key not in st.session_state:
        st.session_state[key] = default

# ── User stats helper ─────────────────────────────────────────────────
def _get_user_stats(username: str):
    stories = get_stories(username)
    total_stories = len(stories)
    total_words   = sum(len(s["story"].split()) for s in stories)
    favourite_genre = max(
        (s["genre"] for s in stories if s["genre"]),
        key=lambda g: sum(1 for s in stories if s["genre"] == g),
        default="Fantasy"
    )
    return total_stories, total_words, favourite_genre

# ── Story filtering function ───────────────────────────────────────────
def get_filtered_stories(username: str, search_query: str = "", genre_filter: str = "All", 
                        start_date=None, end_date=None, time_filter: str = "All Time"):
    """Filter stories based on search criteria"""
    all_stories = get_stories(username)
    
    # Apply text search filter
    if search_query and search_query.strip():
        search_terms = search_query.lower().split()
        all_stories = [
            story for story in all_stories
            if any(term in story['prompt'].lower() for term in search_terms)
        ]
    
    # Apply genre filter
    if genre_filter != "All":
        all_stories = [story for story in all_stories if story['genre'] == genre_filter]
    
    # Apply date range filter
    if start_date or end_date:
        filtered_stories = []
        for story in all_stories:
            try:
                story_date = datetime.strptime(story['created_at'], "%Y-%m-%d %H:%M:%S").date()
                if start_date and story_date < start_date:
                    continue
                if end_date and story_date > end_date:
                    continue
                filtered_stories.append(story)
            except (ValueError, KeyError):
                continue
        all_stories = filtered_stories
    
    # Apply time period filter
    if time_filter != "All Time":
        now = datetime.now()
        filtered_stories = []
        
        for story in all_stories:
            try:
                story_date = datetime.strptime(story['created_at'], "%Y-%m-%d %H:%M:%S")
                
                if time_filter == "Today":
                    if story_date.date() == now.date():
                        filtered_stories.append(story)
                elif time_filter == "This Week":
                    week_ago = now - timedelta(days=7)
                    if story_date >= week_ago:
                        filtered_stories.append(story)
                elif time_filter == "This Month":
                    month_ago = now - timedelta(days=30)
                    if story_date >= month_ago:
                        filtered_stories.append(story)
                elif time_filter == "Last 3 Months":
                    three_months_ago = now - timedelta(days=90)
                    if story_date >= three_months_ago:
                        filtered_stories.append(story)
            except (ValueError, KeyError):
                continue
        
        all_stories = filtered_stories
    
    # Sort by most recent first
    all_stories.sort(key=lambda x: x.get('created_at', ''), reverse=True)
    
    return all_stories

# ── Export helpers ─────────────────────────────────────────────────────
_GENRE_ACCENT = {
    "Mixed":     (  0, 212, 170),
    "Romantic":  (220, 60, 110),
    "Horror":    ( 20, 200, 50),
    "Fantasy":   (130, 80, 220),
    "Mystery":   ( 60, 90, 200),
    "Humor":     (220, 170, 20),
    "Suspense":  (210, 30, 30),
    "Sci-Fi":    ( 20, 130, 230),
    "Adventure": ( 50, 170, 50),
}
_GENRE_BG_RGB = {
    "Mixed":     (10, 15, 25),
    "Romantic":  (45, 10, 30),
    "Horror":    (10, 10, 10),
    "Fantasy":   (16,   8, 48),
    "Mystery":   ( 9, 13, 26),
    "Humor":     (20, 18,   4),
    "Suspense":  (31,   3,   3),
    "Sci-Fi":    ( 2, 10, 31),
    "Adventure": ( 6, 14,   2),
}

def _safe_text(text: str) -> str:
    """Transliterate common Unicode to latin-1 for FPDF built-in fonts."""
    subs = {
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u2013": "-", "\u2014": "--", "\u2026": "...", "\u00a0": " ",
        "\u2022": "*", "\u2032": "'",
    }
    for k, v in subs.items():
        text = text.replace(k, v)
    return text.encode("latin-1", errors="replace").decode("latin-1")

class _StoryPDF(FPDF):
    def __init__(self, accent_rgb, genre):
        super().__init__()
        self._ar, self._ag, self._ab = accent_rgb
        self._genre = genre

    def footer(self):
        self.set_y(-12)
        self.set_draw_color(self._ar, self._ag, self._ab)
        self.line(10, self.get_y(), 200, self.get_y())
        self.set_y(-10)
        self.set_font("Helvetica", "", 7)
        self.set_text_color(150, 150, 160)
        self.cell(95, 5, "Generated by Story Waver")
        self.cell(95, 5, f"Page {self.page_no()}", align="R")

@st.cache_data(show_spinner=False)
def make_pdf(sid, title, genre, tone, mode, prompt, story_text, created_at) -> bytes:
    accent = _GENRE_ACCENT.get(genre, (0, 212, 170))
    r, g, b = accent

    pdf = _StoryPDF(accent, genre)
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()

    # Header bar
    pdf.set_fill_color(r, g, b)
    pdf.rect(0, 0, 210, 14, "F")
    pdf.set_xy(8, 3)
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 8, _safe_text(f"STORY WAVER  |  {genre.upper()}  |  {tone.upper()}  |  {mode.upper()}"))

    # Title
    pdf.set_xy(10, 20)
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(30, 30, 50)
    pdf.multi_cell(190, 10, _safe_text(title))

    # Prompt
    pdf.set_x(10)
    pdf.set_font("Helvetica", "I", 9)
    pdf.set_text_color(110, 110, 140)
    pdf.multi_cell(190, 5, _safe_text(f'Prompt: "{prompt}"'))

    # Date
    pdf.set_x(10)
    pdf.set_font("Helvetica", "", 8)
    pdf.set_text_color(150, 150, 160)
    pdf.cell(0, 6, _safe_text(format_date(created_at)))

    # Divider
    y = pdf.get_y() + 4
    pdf.set_draw_color(r, g, b)
    pdf.set_line_width(0.6)
    pdf.line(10, y, 200, y)

    # Story body
    pdf.set_xy(10, y + 6)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(30, 30, 50)
    pdf.multi_cell(190, 6.5, _safe_text(story_text))

    return bytes(pdf.output())

@st.cache_data(show_spinner=False)
def make_pptx(sid, title, genre, tone, mode, prompt, story_text, created_at) -> bytes:
    ar, ag, ab = _GENRE_ACCENT.get(genre, (0, 212, 170))
    br, bg_, bb = _GENRE_BG_RGB.get(genre, (15, 15, 26))
    accent = RGBColor(ar, ag, ab)
    bg_col = RGBColor(br, bg_, bb)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H  = prs.slide_width, prs.slide_height

    def _bg(slide):
        s = slide.shapes.add_shape(1, 0, 0, W, H)
        s.fill.solid(); s.fill.fore_color.rgb = bg_col; s.line.fill.background()

    def _txt(slide, l, t, w, h, text, size, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = wrap
        p  = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
        run.font.color.rgb = color or RGBColor(210, 214, 230)

    # ── Title slide ──
    sl = prs.slides.add_slide(blank)
    _bg(sl)
    _txt(sl, Inches(0.5), Inches(0.2), Inches(5), Inches(0.4),
         "STORY WAVER", 9, bold=True, color=accent)
    _txt(sl, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         f"{genre}  ·  {tone}  ·  {mode}", 10, color=RGBColor(160, 165, 185))
    _txt(sl, Inches(0.5), Inches(1.6), Inches(2.5), Inches(0.4),
         title, 34, bold=True, wrap=True)
    line = sl.shapes.add_shape(1, Inches(0.5), Inches(4.15), Inches(2.5), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
    _txt(sl, Inches(0.5), Inches(4.45), Inches(12), Inches(0.7),
         f'"{prompt}"', 11, italic=True, color=RGBColor(130, 150, 140))
    _txt(sl, Inches(0.5), Inches(5.2), Inches(6), Inches(0.4),
         format_date(created_at), 9, color=RGBColor(90, 90, 110))

    # ── Story content slides ──
    paras = [p for p in story_text.split("\n") if p.strip()]
    chunks, cur = [], ""
    for para in paras:
        if len(cur) + len(para) + 2 < 650:
            cur += ("\n\n" if cur else "") + para
        else:
            if cur: chunks.append(cur)
            cur = para
    if cur:
        chunks.append(cur)

    total = len(chunks) + 1
    for idx, chunk in enumerate(chunks):
        sl = prs.slides.add_slide(blank)
        _bg(sl)
        _txt(sl, Inches(10.8), Inches(0.12), Inches(2.3), Inches(0.3),
             genre.upper(), 8, bold=True, color=accent,
             align=PP_ALIGN.RIGHT, wrap=False)
        _txt(sl, Inches(0.6), Inches(0.55), Inches(12.1), Inches(6.55),
             chunk, 14, wrap=True)
        _txt(sl, Inches(11.3), Inches(7.1), Inches(1.8), Inches(0.3),
             f"{idx + 2} / {total}", 8,
             color=RGBColor(80, 80, 100), align=PP_ALIGN.RIGHT, wrap=False)

    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ── Genre Background Themes ────────────────────────────────────────────
def apply_genre_background(genre: str):
    """Apply genre-specific background without floating emojis."""
    # Baby pink gradient background
    img_url = "linear-gradient(135deg, #fff5f7 0%, #ffe4ec 25%, #ffd9e6 50%, #ffc0cb 75%, #ffb6c1 100%)"
    
    CONFIGS = {
        "Mixed": {"bg": img_url},
        "Romantic": {"bg": img_url},
        "Horror": {"bg": img_url},
        "Fantasy": {"bg": img_url},
        "Mystery": {"bg": img_url},
        "Humor": {"bg": img_url},
        "Suspense": {"bg": img_url},
        "Sci-Fi": {"bg": img_url},
        "Adventure": {"bg": img_url},
    }

    cfg = CONFIGS.get(genre)
    if not cfg:
        return

    # Apply only the background, no floating emojis
    st.markdown(
        f"<style>.stApp{{background:{cfg['bg']} !important;}}</style>",
        unsafe_allow_html=True,
    )

# ════════════════════════════════════════════════════════════
#  AUTH PAGE - Full Screen Professional Design
# ════════════════════════════════════════════════════════════
def show_auth_page():
    # Remove the auth-fullscreen and auth-content divs that constrain the width
    # Let the content take the full page naturally
    
    st.markdown("""
<style>
/* Override the previous auth styles to make it full page */
.auth-title {
    font-size: 4rem;
    font-weight: 900;
    letter-spacing: 12px;
    background: linear-gradient(135deg, #ff69b4, #ff1493);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 10px;
    line-height: 1.2;
    text-align: center;
    padding-top: 40px;
}
.auth-subtitle {
    font-size: 1.3rem;
    color: #ff1493;
    margin-bottom: 30px;
    font-weight: 500;
    letter-spacing: 2px;
    text-align: center;
}
/* Make tabs container full width */
.stTabs [data-baseweb="tab-list"] {
    gap: 12px;
    background: #fff0f5;
    border-radius: 16px;
    padding: 6px;
    margin-bottom: 30px;
    border: 1px solid #ffb6c1;
    display: flex;
    justify-content: center;
    align-items: center;
    max-width: 600px;
    margin-left: auto;
    margin-right: auto;
}
.stTabs [data-baseweb="tab"] {
    border-radius: 12px;
    padding: 12px 24px;
    color: #1e3a8a !important;
    font-weight: 600;
    font-size: 1rem;
    transition: all 0.3s;
}
.stTabs [data-baseweb="tab"][aria-selected="true"] {
    background: linear-gradient(135deg, #ff69b4, #ff1493);
    color: white !important;
}
/* Style the forms to be centered but not constrained */
.stForm {
    max-width: 400px;
    margin-left: auto;
    margin-right: auto;
    padding: 20px;
    background: white;
    border-radius: 24px;
    box-shadow: 0 20px 60px rgba(255, 105, 180, 0.25);
    border: 1px solid rgba(255, 182, 193, 0.3);
}
.stTextInput label {
    color: #1e3a8a !important;
    font-weight: 600 !important;
    font-size: 0.95rem !important;
    margin-bottom: 8px !important;
    text-align: left !important;
    display: block !important;
}
.stTextInput input {
    background: linear-gradient(145deg, #ffe4ec, #ffd9e6) !important;
    border: 2px solid #ffb6c1 !important;
    border-radius: 12px !important;
    padding: 14px 16px !important;
    color: #1e3a8a !important;
    font-size: 1rem !important;
    transition: all 0.3s !important;
    caret-color: #ff1493 !important;
}
.stTextInput input:focus {
    border-color: #ff69b4 !important;
    box-shadow: 0 0 0 4px rgba(255, 105, 180, 0.15) !important;
    outline: none !important;
}
.stTextInput input::placeholder {
    color: #ffb6c1 !important;
    opacity: 0.8 !important;
}
.stTextInput input:hover {
    border-color: #ff69b4 !important;
}
.stButton > button {
    background: linear-gradient(135deg, #ff69b4, #ff1493) !important;
    color: white !important;
    border: none !important;
    border-radius: 12px !important;
    padding: 14px 24px !important;
    font-weight: 600 !important;
    font-size: 1rem !important;
    transition: all 0.3s !important;
    border: 2px solid transparent !important;
    width: 100%;
    margin-top: 10px;
}
.stButton > button:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 25px rgba(255, 105, 180, 0.4) !important;
    border-color: #ffb6c1 !important;
}
.stAlert {
    border-radius: 12px !important;
    border-left: 4px solid #ff69b4 !important;
    background: #fff5f7 !important;
    color: #1e3a8a !important;
    margin-top: 20px !important;
    max-width: 400px;
    margin-left: auto;
    margin-right: auto;
}
.stAlert > div {
    color: #1e3a8a !important;
}
/* Ensure cursor is visible in all browsers */
input, textarea {
    caret-color: #ff1493 !important;
}
/* Add some breathing room at the bottom */
.main > div {
    padding-bottom: 40px !important;
}
</style>
""", unsafe_allow_html=True)

    # Title and subtitle - full width
    st.markdown("<div class='auth-title'>STORY WAVER</div>", unsafe_allow_html=True)
    st.markdown("<div class='auth-subtitle'>Your Personal AI Story Companion</div>", unsafe_allow_html=True)

    # Inject JS so the browser password manager can autofill saved credentials
    st.components.v1.html("""
    <script>
    (function() {
        function applyAutocomplete() {
            var doc = window.parent.document;
            doc.querySelectorAll('[data-testid="stTextInput"]').forEach(function(container) {
                var label = container.querySelector('label');
                var input = container.querySelector('input');
                if (!label || !input) return;
                var text = label.textContent.toLowerCase().trim();
                if (text === 'username' || text === 'choose username') {
                    input.setAttribute('autocomplete', text === 'choose username' ? 'new-username' : 'username');
                    input.setAttribute('name', text === 'choose username' ? 'new-username' : 'username');
                } else if (text === 'password' || text === 'choose password') {
                    input.setAttribute('autocomplete', text === 'choose password' ? 'new-password' : 'current-password');
                    input.setAttribute('name', text === 'choose password' ? 'new-password' : 'password');
                } else if (text.includes('confirm')) {
                    input.setAttribute('autocomplete', 'new-password');
                    input.setAttribute('name', 'confirm-password');
                }
            });
        }
        applyAutocomplete();
        setInterval(applyAutocomplete, 400);
    })();
    </script>
    """, height=0)

    tab_login, tab_register = st.tabs(["Login", "Register"])

    with tab_login:
        st.markdown("<br>", unsafe_allow_html=True)
        with st.form("login_form"):
            uname = st.text_input("Username", key="login_user", placeholder="Enter your username")
            pwd   = st.text_input("Password", type="password", key="login_pwd", placeholder="•••••••")
            st.markdown("<br>", unsafe_allow_html=True)
            submitted = st.form_submit_button("Login", type="primary", use_container_width=True)
        if submitted:
            if uname and pwd:
                ok, msg = verify_user(uname, pwd)
                if ok:
                    st.session_state.logged_in = True
                    st.session_state.username = uname.strip().lower()
                    st.session_state.messages = []
                    st.rerun()
                else:
                    st.error(msg)
            else:
                st.warning("Please enter username and password.")

    with tab_register:
        st.markdown("<br>", unsafe_allow_html=True)
        with st.form("register_form"):
            new_user = st.text_input("Choose Username", key="reg_user",  placeholder="Minimum 3 characters")
            new_pwd  = st.text_input("Choose Password", type="password", key="reg_pwd",  placeholder="Minimum 4 characters")
            new_pwd2 = st.text_input("Confirm Password", type="password", key="reg_pwd2", placeholder="Repeat password")
            st.markdown("<br>", unsafe_allow_html=True)
            reg_submitted = st.form_submit_button("Create Account", type="primary", use_container_width=True)
        if reg_submitted:
            if new_pwd != new_pwd2:
                st.error("Passwords do not match.")
            elif new_user and new_pwd:
                ok, msg = create_user(new_user, new_pwd)
                if ok:
                    st.success(msg + " Please login.")
                else:
                    st.error(msg)
            else:
                st.warning("Please fill in all fields.")

# ════════════════════════════════════════════════════════════
#  HISTORY PAGE
# ════════════════════════════════════════════════════════════════
def show_history_page():
    username = st.session_state.username
    stories  = get_stories(username)

    st.markdown("<div class='main-title' style='font-size:2rem;'>MY STORIES</div>", unsafe_allow_html=True)
    st.markdown(
        f"<div class='main-subtitle'>{len(stories)} stories saved for "
        f"<b style='color:#ff1493'>@{username}</b></div>",
        unsafe_allow_html=True,
    )

    if st.button("← Back to Writer", key="back_from_history"):
        st.session_state.view_story_id = None
        st.session_state.page = "main"
        st.rerun()

    st.markdown("---")

    # ── Single story view ──
    if st.session_state.view_story_id:
        story = get_story_by_id(username, st.session_state.view_story_id)
        if story:
            # UPDATED: Added 'selected' class to make genre/writing mode/tone appear in white
            st.markdown(
                f"<div class='story-view-header'>"
                f"<span class='tag-genre selected'>{story['genre']}</span>"
                f"<span class='tag-mode selected'>{story['writing_mode']}</span>"
                f"<span class='tag-tone selected'>{story['tone']}</span>"
                f"<div style='color:#ccccff;font-size:0.95rem;font-weight:500;margin-top:10px;line-height:1.4;'>{story['title']}</div>"
                f"<div class='story-view-prompt'><span>Prompt:</span>{story['prompt']}</div>"
                f"<div class='story-view-meta'>{format_date(story['created_at'])}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
            st.markdown(f"<div class='story-box'>{story['story']}</div>", unsafe_allow_html=True)
            c1, c2 = st.columns([3, 2])
            with c1:
                if st.button("← Back to History", key="back_to_list"):
                    st.session_state.view_story_id = None
                    st.rerun()
            with c2:
                if st.button("Load into Writer", type="primary", key="load_story", use_container_width=True):
                    # Inject as a message so it appears in chat
                    st.session_state.messages = [
                        {"role": "user",      "content": story["prompt"]},
                        {"role": "assistant", "content": story["story"],
                         "genre": story["genre"], "writing_mode": story["writing_mode"],
                         "tone": story["tone"]},
                    ]
                    st.session_state.view_story_id = None
                    st.session_state.page = "main"
                    st.rerun()

            # ── Export row ──
            st.markdown(
                "<div style='margin-top:10px;margin-bottom:4px;color:#ff69b4;"
                "font-size:0.75rem;letter-spacing:1.2px;text-transform:uppercase;"
                "font-weight:700;'>⬇ Export Story As</div>",
                unsafe_allow_html=True,
            )
            ex1, ex2, ex3 = st.columns(3)
            sid = story["id"]
            with ex1:
                st.download_button(
                    "📄 PDF",
                    data=make_pdf(sid, story["title"], story["genre"], story["tone"],
                                  story["writing_mode"], story["prompt"],
                                  story["story"], story["created_at"]),
                    file_name=f"story_{sid}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key=f"dl_pdf_{sid}",
                )
            with ex2:
                st.download_button(
                    "📊 PPT",
                    data=make_pptx(sid, story["title"], story["genre"], story["tone"],
                                   story["writing_mode"], story["prompt"],
                                   story["story"], story["created_at"]),
                    file_name=f"story_{sid}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    key=f"dl_ppt_{sid}",
                )
            with ex3:
                st.download_button(
                    "📝 TXT",
                    data=story["story"],
                    file_name=f"story_{sid}.txt",
                    mime="text/plain",
                    use_container_width=True,
                    key=f"dl_txt_{sid}",
                )
        return

    # ── Story list ──
    if not stories:
        st.markdown(
            "<div style='text-align:center;color:#ff69b4;padding:60px 0;font-size:1.1rem;'>"
            "📜 No stories yet — go write your first one!</div>",
            unsafe_allow_html=True,
        )
        return

    all_genres   = sorted(set(s["genre"] for s in stories))
    genre_filter = st.selectbox("Filter by Genre", ["All"] + all_genres, key="hist_filter")
    filtered     = stories if genre_filter == "All" else [s for s in stories if s["genre"] == genre_filter]

    st.markdown(f"<div style='color:#ff69b4;font-size:0.8rem;margin-bottom:12px;'>{len(filtered)} stories</div>",
                unsafe_allow_html=True)

    for story in filtered:
        c1, c2 = st.columns([6, 1])
        with c1:
            _prompt_snippet = story['prompt'][:90] + "…" if len(story['prompt']) > 90 else story['prompt']
            # UPDATED: Added 'selected' class to genre tag
            st.markdown(
                f"<div class='history-card'>"
                f"<span class='tag-genre selected'>{story['genre']}</span>"
                f"<div class='history-prompt-label'>Prompt</div>"
                f"<div class='history-title'>{_prompt_snippet}</div>"
                f"<div class='history-date'>{format_date(story['created_at'])}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
        with c2:
            st.markdown("<div style='padding-top:8px;'>", unsafe_allow_html=True)
            if st.button("View",   key=f"view_{story['id']}", use_container_width=True):
                st.session_state.view_story_id = story["id"]
                st.rerun()
            if st.button("Delete", key=f"del_{story['id']}",  use_container_width=True):
                delete_story(username, story["id"])
                st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)

# ════════════════════════════════════════════════════════════════
#  MAIN APP (CHAT INTERFACE)
# ════════════════════════════════════════════════════════════════════
def show_main_app():
    if not st.session_state.get("username"):
        st.error("Please log in to continue.")
        return
    
    username     = st.session_state.username
    info         = get_user_info(username)
    is_connected = check_ollama_connection()
    model_name   = get_model_name()

    # Grab any pending generation queued by sidebar buttons BEFORE rendering sidebar
    pending_prompt = st.session_state.pop("pending_prompt", None)
    pending_mode   = st.session_state.pop("pending_mode",   "generate")

    # Create layout with bottom sidebar
    left_sidebar, main_content, right_sidebar = st.columns([1.5, 3, 1.2])
    
    # Create bottom sidebar
    bottom_sidebar = st.container()

    # ── Left Sidebar ───────────────────────────────────────────────────
    with left_sidebar:
        try:
            with open("logo.png", "rb") as _f:
                _logo_b64 = base64.b64encode(_f.read()).decode()
            _logo_html = f"<img src='data:image/png;base64,{_logo_b64}' style='width:180px;margin-bottom:4px;'>"
        except FileNotFoundError:
            try:
                with open("logo.jpeg", "rb") as _f:
                    _logo_b64 = base64.b64encode(_f.read()).decode()
                _logo_html = f"<img src='data:image/jpeg;base64,{_logo_b64}' style='width:180px;margin-bottom:4px;'>"
            except FileNotFoundError:
                _logo_html = "<span style='font-size:1.8rem;'>📖</span>"
        st.markdown(
            f"<div style='text-align:center;padding:10px 0 4px;'>"
            f"{_logo_html}"
            f"</div>",
            unsafe_allow_html=True,
        )
        st.markdown(
            f"<div style='text-align:center;margin:6px 0;'>"
            f"<span class='user-badge'>👤 {username}</span></div>",
            unsafe_allow_html=True,
        )
        if is_connected:
            st.markdown(
                f"<div style='text-align:center;margin-bottom:4px;'>"
                f"<span class='status-online'>● ONLINE</span></div>",
                unsafe_allow_html=True,
            )
        else:
            st.markdown(
                "<div style='text-align:center;margin-bottom:4px;'>"
                "<span class='status-offline'>● OFFLINE — Start Ollama</span></div>",
                unsafe_allow_html=True,
            )

        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Story Controls</div>", unsafe_allow_html=True)

        genre          = st.selectbox("Genre",  ["Romantic","Horror","Fantasy","Mystery","Humor","Suspense","Sci-Fi","Adventure"])
        writing_mode   = st.selectbox("Writing Mode", ["Beginning","Continue","Climax","Ending"])
        tone           = st.selectbox("Tone",   ["Dramatic","Lighthearted","Dark","Whimsical","Suspenseful","Melancholic","Epic"])
        creativity     = st.slider("Creativity", 0, 100, 50)
        response_length = st.selectbox("Response Length", ["Short","Medium","Long"], index=1)

        # Interactive Plot Steering Section
        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Interactive Plot Steering</div>", unsafe_allow_html=True)
        
        plot_steering_options = [
            "Introduce a twist",
            "Reveal a secret", 
            "Add a new character",
            "Increase conflict",
            "Develop romance",
            "Create mystery",
            "Build suspense",
            "Flashback"
        ]
        
        plot_steering = st.selectbox(
            "Choose plot direction (optional)",
            ["None"] + plot_steering_options,
            key="plot_steering",
            help="Select how you want story to develop in next paragraph"
        )

        # Plot Steering Start Button (appears right below dropdown)
        if plot_steering != "None" and any(m["role"] == "assistant" for m in st.session_state.messages):
            if st.button(f"🎬 Start {plot_steering}", type="primary", use_container_width=True, key="plot_steering_start"):
                # Store the plot steering action in session state to be processed later
                st.session_state.pending_plot_steering = plot_steering
                st.rerun()

        
        has_story = any(m["role"] == "assistant" for m in st.session_state.messages)

        # Collaborative Story Mode Section
        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Collaborative Story Mode</div>", unsafe_allow_html=True)
        
        collaborative_mode = st.checkbox(
            "Enable Collaborative Mode",
            value=st.session_state.collaborative_mode,
            key="collaborative_mode_toggle",
            help="Multiple users take turns contributing to story, with AI filling gaps"
        )
        
        # Update session state when checkbox changes
        if collaborative_mode != st.session_state.collaborative_mode:
            st.session_state.collaborative_mode = collaborative_mode
            if collaborative_mode:
                st.session_state.collaborative_turn = 1
                st.session_state.collaborative_contributions = []
            st.rerun()
        
        if st.session_state.collaborative_mode:
            st.markdown(f"""
            <div class='collaborative-indicator'>
                <div class='collaborative-turn'>🤝 Current Turn</div>
                <div class='collaborative-user'>User {st.session_state.collaborative_turn}'s turn to write</div>
                <div class='collaborative-info'>AI will fill gaps between contributions</div>
            </div>
            """, unsafe_allow_html=True)

        if st.button("Continue Story", type="primary", use_container_width=True,
                     disabled=not is_connected or not has_story):
            st.session_state.pending_prompt = "Continue story in an engaging and creative way."
            st.session_state.pending_mode   = "continue"
            st.rerun()

    # ── Right Sidebar ──────────────────────────────────────────────────
    with right_sidebar:
        st.markdown("<div class='sidebar-section-header'>📚 Story History</div>", unsafe_allow_html=True)
        
        # Search and Filter Controls
        with st.expander("🔍 Search & Filter", expanded=False):
            # Text search
            search_query = st.text_input("Search by prompt:", placeholder="Enter keywords...", key="history_search")
            
            # Genre filter
            all_genres = ["All", "Fantasy", "Mystery", "Sci-Fi", "Romance", "Horror", "Adventure", "Historical", "Comedy", "Drama"]
            selected_genre = st.selectbox("Filter by genre:", all_genres, key="genre_filter")
            
            # Date range filter (Removed per user request)
            start_date = None
            end_date = None
            
            # Time filter
            time_filter = st.selectbox("Time period:", ["All Time", "Today", "This Week", "This Month", "Last 3 Months"], key="time_filter")
            
            # Clear filters button
            if st.button("🗑️ Clear Filters", use_container_width=True):
                st.session_state.history_search = ""
                st.session_state.genre_filter = "All"
                st.session_state.start_date = None
                st.session_state.end_date = None
                st.session_state.time_filter = "All Time"
                st.rerun()
        
        # Apply filters to get stories
        filtered_stories = get_filtered_stories(username, search_query, selected_genre, start_date, end_date, time_filter)
        
        # Display filtered results
        if filtered_stories:
            st.markdown(f"<div style='color:#ff69b4;font-size:0.75rem;margin-bottom:8px;'>{len(filtered_stories)} stories found</div>", unsafe_allow_html=True)
            
            for s in filtered_stories[:5]:  # Show up to 5 results
                _s_snippet = s['prompt'][:70] + "…" if len(s['prompt']) > 70 else s['prompt']
                # UPDATED: Added 'selected' class to genre tag
                st.markdown(
                    f"<div class='history-card'>"
                    f"<span class='tag-genre selected'>{s['genre']}</span>"
                    f"<div class='history-prompt-label'>Prompt</div>"
                    f"<div class='history-title'>{_s_snippet}</div>"
                    f"<div class='history-date'>{format_date(s['created_at'])}</div>"
                    f"</div>",
                    unsafe_allow_html=True,
                )
            
            if len(filtered_stories) > 5:
                if st.button(f"View All {len(filtered_stories)} Stories →", use_container_width=True, key="goto_history"):
                    st.session_state.page = "history"
                    st.rerun()
        else:
            st.markdown(
                "<div style='color:#ff69b4;font-size:0.82rem;text-align:center;'>No stories found matching your filters!</div>",
                unsafe_allow_html=True,
            )

        # Add New Story button below chat history
        st.markdown("---")
        if st.button("📝 New Story", type="primary", use_container_width=True):
            st.session_state.messages = []
            st.rerun()
        
        # Add Logout button below New Story
        if st.button("🚪 Logout", use_container_width=True):
            st.session_state.logged_in = False
            st.session_state.username  = ""
            st.session_state.messages  = []
            st.session_state.page      = "main"
            st.rerun()

    # ── Main Content Area ───────────────────────────────────────────────
    with main_content:
        # Apply genre-specific animated background theme
        apply_genre_background(genre)

        # ── Main chat area ────────────────────────────────────────────────
        st.markdown("<div class='main-title'>STORY WAVER</div>", unsafe_allow_html=True)
        st.markdown("<div class='main-subtitle'>Begin Your Creative Story Journey ✨</div>",
                    unsafe_allow_html=True)

    # ── Stats dashboard (always shown on welcome) ──
        if not st.session_state.messages and not pending_prompt:
            _total, _words, _fav = _get_user_stats(username)
            _words_str = f"{_words:,}" if _words >= 1000 else str(_words)
            st.markdown(
                f"<div class='stats-row'>"
                f"<div class='stat-card'><span class='stat-icon'>📚</span>"
                f"<span class='stat-value'>{_total}</span>"
                f"<span class='stat-label'>Stories Created</span></div>"
                f"<div class='stat-card'><span class='stat-icon'>✍️</span>"
                f"<span class='stat-value'>{_words_str}</span>"
                f"<span class='stat-label'>Words Written</span></div>"
                f"<div class='stat-card'><span class='stat-icon'>🎭</span>"
                f"<span class='stat-value' style='font-size:1.1rem;letter-spacing:-0.5px;'>{_fav}</span>"
                f"<span class='stat-label'>Favourite Genre</span></div>"
                f"</div>",
                unsafe_allow_html=True,
            )

        
        # ── Story phase progress bar (always visible) ──
        _phases      = ["Beginning", "Continue", "Climax", "Ending"]
        _phase_idx   = _phases.index(writing_mode) if writing_mode in _phases else 0
        _session_wc  = sum(len(m["content"].split()) for m in st.session_state.messages
                           if m["role"] == "assistant" and not m.get("blocked"))
        _nodes_html  = ""
        for idx, phase in enumerate(_phases):
            if idx < _phase_idx:
                dot_cls = "done"; lbl_cls = "done"
            elif idx == _phase_idx:
                dot_cls = "active"; lbl_cls = "active"
            else:
                dot_cls = ""; lbl_cls = ""
            _nodes_html += f"<div class='phase-node'><div class='phase-dot {dot_cls}'></div><div class='phase-label {lbl_cls}'>{phase}</div></div>"
            if idx < len(_phases) - 1:
                line_cls = "done" if idx < _phase_idx else ""
                _nodes_html += f"<div class='phase-line {line_cls}'></div>"
        _words_note = f"✍️ {_session_wc:,} words in current session" if _session_wc > 0 else "✍️ Start writing to track your word count"
        st.markdown(
            f"<div class='progress-wrap'>"
            f"<div class='progress-title'>Story Progress</div>"
            f"<div class='progress-phases'>{_nodes_html}</div>"
            f"<div class='progress-words'>{_words_note}</div>"
            f"</div>",
            unsafe_allow_html=True,
        )

        # ── Render existing conversation ──
        for msg in st.session_state.messages:
            if msg["role"] == "user":
                # Handle collaborative user messages
                if msg.get("collaborative_turn"):
                    with st.chat_message("user"):
                        st.markdown(msg["content"])
                else:
                    with st.chat_message("user"):
                        st.markdown(msg["content"])
            else:
                with st.chat_message("assistant"):
                    if msg.get("blocked"):
                        st.markdown(msg["content"], unsafe_allow_html=True)
                    else:
                        # Handle collaborative AI messages
                        if msg.get("collaborative"):
                            st.markdown("**🤖 AI Gap Filler:**")
                        else:
                            # UPDATED: Add 'selected' class to tags that match current selections
                            current_genre = genre if 'genre' in locals() else msg.get('genre', '')
                            genre_class = "tag-genre selected" if msg.get('genre') == current_genre else "tag-genre"
                            
                            current_mode = writing_mode if 'writing_mode' in locals() else msg.get('writing_mode', '')
                            mode_class = "tag-mode selected" if msg.get('writing_mode') == current_mode else "tag-mode"
                            
                            current_tone = tone if 'tone' in locals() else msg.get('tone', '')
                            tone_class = "tag-tone selected" if msg.get('tone') == current_tone else "tag-tone"
                            
                            tags_html = (
                                f"<span class='{genre_class}'>{msg.get('genre','')}</span>"
                                f"<span class='{mode_class}'>{msg.get('writing_mode','')}</span>"
                                f"<span class='{tone_class}'>{msg.get('tone','')}</span>"
                            )
                            if msg.get("plot_steering"):
                                tags_html += f"<span class='msg-tag tag-mode' style='background:#ffb6c1;color:#1e3a8a;border-color:#ff69b4;'>🎭 {msg.get('plot_steering')}</span>"
                            
                            st.markdown(tags_html, unsafe_allow_html=True)
                        st.markdown(msg["content"])

    # ── Generation logic ──────────────────────────────────────────────
    def run_generation(prompt_text: str, mode: str = "generate", plot_steering: str = "None"):
        """Handle regular story generation."""
        if not prompt_text or not prompt_text.strip():
            return False
            
        effective_input = prompt_text.strip()
        effective_mode = mode if mode != "generate" else (
            "Continue" if st.session_state.messages else "Beginning"
        )
        
        # Show user message
        with st.chat_message("user"):
            st.markdown(effective_input)
        
        # Generate AI response
        context = ""
        if st.session_state.messages:
            for m in reversed(st.session_state.messages):
                if m["role"] == "assistant":
                    context = m["content"]
                    break
        
        effective_plot_steering = plot_steering if plot_steering != "None" else ""
        
        # Build character information from session state
        character_info = ""
        if st.session_state.get("story_characters"):
            chars = st.session_state.story_characters
            if chars.get("protagonist"):
                character_info += f"Main Character: {chars['protagonist']}"
                if chars.get("protagonist_role"):
                    character_info += f" ({chars['protagonist_role']})"
                character_info += ". "
            if chars.get("antagonist"):
                character_info += f"Antagonist: {chars['antagonist']}"
                if chars.get("antagonist_role"):
                    character_info += f" ({chars['antagonist_role']})"
                character_info += ". "
        
        # Add story elements and theme
        elements_info = ""
        if st.session_state.get("selected_story_elements"):
            elements_info = f"Story Elements: {', '.join(st.session_state.selected_story_elements)}. "
        if st.session_state.get("selected_story_theme"):
            theme_info = f"Theme: {st.session_state.selected_story_theme}. "
        
        # Combine all additional context
        additional_context = character_info + elements_info + theme_info
        
        with st.chat_message("assistant"):
            full_text = st.write_stream(
                generate_story_stream(
                    user_input=effective_input,
                    genre=genre,
                    writing_mode=effective_mode,
                    tone=tone,
                    creativity=creativity,
                    response_length=response_length,
                    story_context=context,
                    plot_steering=effective_plot_steering,
                    character_info=additional_context,
                )
            )
        
        # Update message history
        st.session_state.messages.append({
            "role": "user", 
            "content": effective_input
        })
        
        message_data = {
            "role": "assistant", "content": full_text,
            "genre": genre, "writing_mode": effective_mode, "tone": tone,
        }
        if effective_plot_steering:
            message_data["plot_steering"] = effective_plot_steering

        st.session_state.messages.append(message_data)
        
        # Save to database
        _clean_prompt = effective_input.strip().strip("?.!,;:-_/\\\"'")
        if full_text.strip() and len(_clean_prompt) >= 8:
            save_story(username=username, prompt=effective_input, story=full_text,
                       genre=genre, writing_mode=effective_mode, tone=tone)

        st.rerun()

    # Process any pending plot steering action
    pending_plot_steering = st.session_state.pop("pending_plot_steering", None)
    if pending_plot_steering:
        # Continue existing story with plot steering
        context = ""
        if st.session_state.messages:
            for m in reversed(st.session_state.messages):
                if m["role"] == "assistant":
                    context = m["content"]
                    break
        
        plot_prompt = f"Continue the story with {pending_plot_steering.lower()}"
        run_generation(plot_prompt, "continue", pending_plot_steering)

    def run_collaborative_generation(prompt_text: str):
            """Handle collaborative story generation with turn-based user contributions."""
            if not st.session_state.collaborative_mode:
                return False
                
            # Build character information from session state
            character_info = ""
            if st.session_state.get("story_characters"):
                chars = st.session_state.story_characters
                if chars.get("protagonist"):
                    character_info += f"Main Character: {chars['protagonist']}"
                    if chars.get("protagonist_role"):
                        character_info += f" ({chars['protagonist_role']})"
                    character_info += ". "
                if chars.get("antagonist"):
                    character_info += f"Antagonist: {chars['antagonist']}"
                    if chars.get("antagonist_role"):
                        character_info += f" ({chars['antagonist_role']})"
                    character_info += ". "
                if chars.get("supporting"):
                    supporting = [c for c in chars["supporting"] if c.strip()]
                    if supporting:
                        character_info += f"Supporting Characters: {', '.join(supporting)}. "
            
            # Add story elements and theme
            elements_info = ""
            if st.session_state.get("selected_story_elements"):
                elements_info = f"Story Elements: {', '.join(st.session_state.selected_story_elements)}. "
            if st.session_state.get("selected_story_theme"):
                theme_info = f"Theme: {st.session_state.selected_story_theme}. "
            
            # Combine all additional context
            additional_context = character_info + elements_info + theme_info
                
            # Store user contribution
            st.session_state.collaborative_contributions.append({
                "user": st.session_state.collaborative_turn,
                "content": prompt_text,
                "timestamp": datetime.now().isoformat()
            })
            
            # Show user contribution
            with st.chat_message("user"):
                st.markdown(f"**👤 User {st.session_state.collaborative_turn}:** {prompt_text}")
            
            # Generate AI gap-filler
            context = ""
            if st.session_state.messages:
                for m in reversed(st.session_state.messages):
                    if m["role"] == "assistant":
                        context = m["content"]
                        break
            
            effective_plot_steering = plot_steering if plot_steering != "None" else ""
            
            with st.chat_message("assistant"):
                st.markdown("**🤖 AI Gap Filler:**")
                full_text = st.write_stream(
                    generate_story_stream(
                        user_input=f"Bridge between User {st.session_state.collaborative_turn}'s contribution and the next part",
                        genre=genre,
                        writing_mode="Continue",
                        tone=tone,
                        creativity=creativity,
                        response_length="Short",
                        story_context=context,
                        plot_steering=effective_plot_steering,
                        character_info=additional_context,
                    )
                )
            
            # Update message history
            st.session_state.messages.append({
                "role": "user", 
                "content": f"User {st.session_state.collaborative_turn}: {prompt_text}",
                "collaborative_turn": st.session_state.collaborative_turn
            })
            
            st.session_state.messages.append({
                "role": "assistant", 
                "content": full_text,
                "genre": genre, 
                "writing_mode": "Continue", 
                "tone": tone,
                "collaborative": True
            })
            
            # Toggle turn
            st.session_state.collaborative_turn = 2 if st.session_state.collaborative_turn == 1 else 1
            
            st.rerun()

    # ── Start Story Generation Button (always shown when characters available) ──
    # Check if character information is available
    has_chars = st.session_state.get("story_characters", {})
    protagonist = has_chars.get("protagonist", "").strip()
    antagonist = has_chars.get("antagonist", "").strip()
    
    if protagonist or antagonist:
        if not st.session_state.messages and not pending_prompt:
            st.markdown("---")
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            if st.button("🚀 Start Story with Characters", type="primary", use_container_width=True):
                # Generate a story based on character information
                character_prompt = f"Create a story featuring"
                if protagonist:
                    character_prompt += f" {protagonist}"
                    if has_chars.get("protagonist_role"):
                        character_prompt += f" as {has_chars['protagonist_role']}"
                if antagonist:
                    if protagonist:
                        character_prompt += f" and {antagonist}"
                    else:
                        character_prompt += f" {antagonist}"
                    if has_chars.get("antagonist_role"):
                        character_prompt += f" as {has_chars['antagonist_role']}"
                character_prompt += "."
                
                run_generation(character_prompt, "generate", plot_steering)

    # ── Chat input — auto-docked to bottom by Streamlit ──
    if st.session_state.collaborative_mode:
        placeholder = f"User {st.session_state.collaborative_turn}'s turn - Add to the story…" if is_connected else "Ollama is offline — run: ollama serve"
    else:
        placeholder = "Describe your story idea…" if is_connected else "Ollama is offline — run: ollama serve"
            
    if prompt := st.chat_input(placeholder, disabled=not is_connected):
        if st.session_state.collaborative_mode:
            run_collaborative_generation(prompt)
        else:
            run_generation(prompt, "generate", plot_steering)

    # ── Bottom Sidebar ────────────────────────────────────────────────
    with bottom_sidebar:
        st.markdown("<div class='sidebar-section-header'>Story Elements</div>", unsafe_allow_html=True)
        
        # Story Elements Section
        story_elements = {
            "🏰": "Castle/Fortress",
            "⚔️": "Sword/Weapon", 
            "🗝️": "Map/Quest",
            "📜": "Ancient Text/Prophecy",
            "🔮": "Magic/Mystery",
            "👑": "Royal/Court",
            "🌊": "Sea/Voyage",
            "🏔": "Forest/Nature",
            "🌙": "Night/Moon",
            "⚡": "Storm/Power"
        }
        
        selected_elements = []
        cols = st.columns(5)
        for i, (emoji, description) in enumerate(story_elements.items()):
            with cols[i % 5]:
                if st.checkbox(f"{emoji} {description}", key=f"element_{i}"):
                    selected_elements.append(f"{emoji} {description}")
        
        if selected_elements:
            st.session_state.selected_story_elements = selected_elements

        # Story Themes Section
        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Story Themes</div>", unsafe_allow_html=True)
        
        story_themes = {
            "🏰": "Castle/Fortress",
            "🌙": "Night/Moon",
            "🔮": "Magic/Mystery",
            "🌊": "Sea/Voyage",
            "🏔": "Forest/Nature",
            "⚔️": "Sword/Weapon",
            "📜": "Ancient Text/Prophecy",
            "👑": "Royal/Court",
            "⚡": "Storm/Power"
        }
        
        selected_theme = st.selectbox(
            "Choose Theme:",
            options=list(story_themes.values()),
            format_func=lambda x: f"🎭 {x}",
            key="story_theme"
        )
        
        if selected_theme:
            st.session_state.selected_story_theme = selected_theme

        # Characters Section
        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Story Characters</div>", unsafe_allow_html=True)
        
        # Character input fields
        with st.expander("Add Characters (Optional)"):
            col1, col2 = st.columns(2)
            with col1:
                protagonist = st.text_input("Protagonist:", placeholder="Main character name", key="protagonist")
                protagonist_role = st.text_input("Role:", placeholder="Hero, detective, student...", key="protagonist_role")
            with col2:
                antagonist = st.text_input("Antagonist:", placeholder="Villain, rival, obstacle...", key="antagonist")
                antagonist_role = st.text_input("Role:", placeholder="Villain, competitor, challenge...", key="antagonist_role")
        
        # Store characters in session state
        characters = {
            "protagonist": protagonist,
            "protagonist_role": protagonist_role,
            "antagonist": antagonist, 
            "antagonist_role": antagonist_role,
        }
        
        if any(characters.values()):
            st.session_state.story_characters = characters

# ══════════════════════════════════════════════════════════════
#  ROUTER
# ══════════════════════════════════════════════════════════════════
if not st.session_state.logged_in:
    show_auth_page()
elif st.session_state.page == "history":
    show_history_page()
else:
    show_main_app()
